"""Build the COMS 6998 final-project slide deck for the super-linear-synergy
paper. Output: paper_slides/synergy_deck.pptx."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ---- Theme ----------------------------------------------------------------
COLUMBIA_BLUE = RGBColor(0x0F, 0x2A, 0x57)        # deep navy
ACCENT_BLUE = RGBColor(0x1F, 0x4E, 0x90)
TEXT_DARK = RGBColor(0x23, 0x2A, 0x33)
TEXT_MUTED = RGBColor(0x5C, 0x66, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT_GREEN = RGBColor(0x1F, 0x8E, 0x4F)
HIGHLIGHT_RED = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, *, text=None, size=None, bold=None, italic=None, color=None, font="Calibri"):
    if text is not None:
        run.text = text
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, *, fill=None, line=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is not None:
        box.line.color.rgb = line
        box.line.width = Pt(1.0)
    else:
        box.line.fill.background()
    return box


def header_bar(slide, title_text):
    """Top header band + slide title."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLUMBIA_BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False

    title = add_textbox(slide, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.65))
    p = title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_run(p.add_run(), text=title_text, size=28, bold=True, color=WHITE)


def footer(slide, page_no, total):
    foot = add_textbox(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    p = foot.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_run(p.add_run(), text="COMS 6998   ·   Jeremiah Mao, Mateo Juliani", size=10, color=TEXT_MUTED)
    p2 = foot.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    # right-aligned page number on a separate paragraph would land below; instead use a 2nd box
    pno = add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.35))
    pp = pno.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    set_run(pp.add_run(), text=f"{page_no} / {total}", size=10, color=TEXT_MUTED)


def add_bullet_list(slide, items, *, left, top, width, height,
                    body_size=20, bullet_color=ACCENT_BLUE, line_spacing=1.18):
    """`items` may be strings (level-0 bullet) or (text, level) tuples."""
    box = add_textbox(slide, left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = level
        para.line_spacing = line_spacing
        para.space_after = Pt(6)
        bullet = "•" if level == 0 else "–"
        set_run(para.add_run(), text=f"{bullet}  ", size=body_size, color=bullet_color, bold=True)
        set_run(para.add_run(), text=text, size=body_size, color=TEXT_DARK)
    return box


def add_table(slide, *, left, top, rows, cols, col_widths, row_heights, header=True):
    width = sum(col_widths, Emu(0))
    height = sum(row_heights, Emu(0))
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    for c, w in enumerate(col_widths):
        tbl.columns[c].width = w
    for r, h in enumerate(row_heights):
        tbl.rows[r].height = h
    return tbl


def fill_cell(cell, text, *, size=14, bold=False, color=TEXT_DARK, fill=None,
              align=PP_ALIGN.CENTER, italic=False):
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    set_run(p.add_run(), text=text, size=size, bold=bold, color=color, italic=italic)


# ---- Slide builders -------------------------------------------------------
def make_blank(prs):
    layout = prs.slide_layouts[6]  # truly blank
    return prs.slides.add_slide(layout)


def slide_title(prs, total):
    s = make_blank(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SLIDE_W, Inches(2.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLUMBIA_BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False

    title_box = add_textbox(s, Inches(0.6), Inches(2.7), Inches(12), Inches(2.2))
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_run(p.add_run(),
            text="A Super-Linear Synergy Between Format Augmentation",
            size=36, bold=True, color=WHITE)
    p2 = title_box.text_frame.add_paragraph()
    set_run(p2.add_run(),
            text="and KL Preservation in Continual LoRA Knowledge Injection",
            size=36, bold=True, color=WHITE)

    sub = add_textbox(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.6))
    p = sub.text_frame.paragraphs[0]
    set_run(p.add_run(), text="Knowledge Integration in Task-Tuned Financial LLMs",
            size=22, color=TEXT_DARK, italic=True)

    auth = add_textbox(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.6))
    p = auth.text_frame.paragraphs[0]
    set_run(p.add_run(), text="Jeremiah Mao   ·   Mateo Juliani",
            size=18, color=TEXT_DARK)

    inst = add_textbox(s, Inches(0.6), Inches(6.45), Inches(12), Inches(0.5))
    p = inst.text_frame.paragraphs[0]
    set_run(p.add_run(), text="COMS 6998   ·   Columbia Engineering",
            size=14, color=TEXT_MUTED)


def slide_motivation(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "Motivation")
    add_bullet_list(s, [
        "Models go stale post-deployment. Knowledge cutoffs are fixed; the world is not.",
        "News-driven domains (finance, current events) require continuous fact updates.",
        "Deployed models are typically task-tuned (e.g., for query decomposition, classification, code, RAG).",
        "Naïve fine-tuning to ingest new facts breaks the underlying task. Editing methods can preserve task ability but often fail to make new facts available across formats.",
        "Goal: a continual update procedure that is cheap, behavioral, and preserves the task while injecting facts that surface across prompt formats.",
    ], left=Inches(0.6), top=Inches(1.2), width=Inches(12.1), height=Inches(5.5),
       body_size=20)
    footer(s, page_no, total)


def slide_qd_task(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "The Task: Financial Query Decomposition")

    intro = add_textbox(s, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.6))
    p = intro.text_frame.paragraphs[0]
    set_run(p.add_run(),
            text='User question →  the model emits 2-4 retrieval sub-queries  →  '
                 'each sub-query is run against a document index.',
            size=18, color=TEXT_DARK)

    q_box = add_textbox(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.55),
                        fill=RGBColor(0xEC, 0xF1, 0xF8))
    pq = q_box.text_frame.paragraphs[0]
    set_run(pq.add_run(), text='User question:  ', size=16, bold=True, color=ACCENT_BLUE)
    set_run(pq.add_run(), text='"How is MCD performing vs competitors?"', size=16, color=TEXT_DARK)

    # Two columns: 2021 vs 2026
    pre_box = add_textbox(s, Inches(0.6), Inches(2.5), Inches(5.95), Inches(4.2),
                          fill=RGBColor(0xF7, 0xF9, 0xFC), line=RGBColor(0xCF, 0xD8, 0xE3))
    p = pre_box.text_frame.paragraphs[0]
    set_run(p.add_run(), text="Pre-2022 (training cutoff)", size=18, bold=True, color=ACCENT_BLUE)
    for sq in [
        '"MCD same-store sales growth vs Wendy\'s 2021"',
        '"Impact of 2021 restaurant labor shortages on MCD operating margins vs competitors"',
        '"Growth of delivery app and drive-thru sales for MCD in 2021"',
        '"Performance of the MCD BTS meal and celebrity partnerships vs competitor marketing"',
    ]:
        para = pre_box.text_frame.add_paragraph()
        para.space_before = Pt(4)
        set_run(para.add_run(), text="•  ", size=14, bold=True, color=ACCENT_BLUE)
        set_run(para.add_run(), text=sq, size=14, color=TEXT_DARK)

    post_box = add_textbox(s, Inches(6.75), Inches(2.5), Inches(5.95), Inches(4.2),
                           fill=RGBColor(0xFD, 0xF7, 0xEF), line=RGBColor(0xE8, 0xC8, 0x9C))
    p = post_box.text_frame.paragraphs[0]
    set_run(p.add_run(), text="Post-2022 (new facts the model needs to learn)",
            size=18, bold=True, color=HIGHLIGHT_RED)
    for sq in [
        '"MCD same-store sales growth vs Wendy\'s 2026"',
        '"ROI and foot traffic impact of MCD\'s 2026 value meal promotions vs competitor discounts"',
        '"Cost savings from MCD\'s AI drive-thru ordering and kitchen automation vs competitors 2026"',
        '"Impact of GLP-1 weight-loss drugs on MCD and fast-food sales volumes 2026"',
    ]:
        para = post_box.text_frame.add_paragraph()
        para.space_before = Pt(4)
        set_run(para.add_run(), text="•  ", size=14, bold=True, color=HIGHLIGHT_RED)
        set_run(para.add_run(), text=sq, size=14, color=TEXT_DARK)

    footer(s, page_no, total)


def slide_gap(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "The Absorption-Integration Gap")
    add_bullet_list(s, [
        "Naïve fine-tuning on the new fact in QA format → the model can produce the answer when asked in QA format …",
        "… but fails to surface the same fact under a QD prompt or other downstream formats. Facts are stored as format-coupled token patterns, not as entity-anchored knowledge.",
        ("Format gap = | F1(QA) − F1(QD) |  on the same probe questions wrapped two ways.", 0),
        ("Across single-environment update methods, the gap sits at 0.07–0.14 F1 — every method that meaningfully absorbs leaves a gap of this magnitude.", 0),
        ("Closing the gap is the central question of this work.", 0),
    ], left=Inches(0.6), top=Inches(1.15), width=Inches(12.1), height=Inches(5.7),
       body_size=20)
    footer(s, page_no, total)


def slide_setup(prs, page_no, total):
    """Setup: facts only. Backbone, update layer, data, evaluation framing.
    Method-side detail (chain, loss, conditions) belongs on the Method slide."""
    s = make_blank(prs)
    header_bar(s, "Setup")
    add_bullet_list(s, [
        ("Backbone:  Qwen3-4B-Instruct-2507, task-tuned via LoRA (r=32 / α=64) on FNSPID financial news with a pre-2022 cutoff. The base model is the QD decomposer from slide 3.", 0),
        ("Update layer:  LoRA r=16 / α=32 on attention Q/K/V/O + MLP up/down/gate.   AdamW, 3 epochs/round, single A10G 24 GB.", 0),
        ("Industrial pattern:  LoRA + LoRA (LoRA-tuned base + LoRA continual updates) is the dominant customer-facing fine-tuning primitive on AWS Bedrock, Together AI, Fireworks, and Vertex AI's open-model garden — the setting matches production practice, not an academic special case.", 0),
        ("Data isolation:  (i) task-tuning corpus (pre-2022 FNSPID) and edit corpus (post-cutoff FNSPID) are temporally separate;  (ii) 96,897 post-cutoff facts are partitioned into 15 disjoint rounds × 200 facts;  (iii) preservation R@10 uses a held-out QD test split (n=104) never seen during task tuning or editing;  (iv) K=5 templates are leak-free — gold answer never appears in any user/system prompt, only in the assistant target.", 0),
        ("No-update baseline confirms edit facts are novel:  before any editing, Qwen3-4B-Instruct-2507 scores ~0.04 abs F1 on the eval probes — essentially zero. So a round-15 abs F1 of 0.411 is real lift from the editing procedure, not pretrained recall. (Caveat: we only fully isolate from the task-tuning corpus, not from the base model's web pretraining.)", 0),
        ("Headline metric (abs F1):  token-F1 of greedy generation vs gold object, averaged over 200 facts × 5 paraphrased probes/round. Two guardrails (preservation R@10, locality F1) on slide 10. Across the COPR family hidden-state geometry does NOT predict behavioral F1 — so we report only behavioral numbers.", 0),
    ], left=Inches(0.6), top=Inches(1.05), width=Inches(12.1), height=Inches(5.85),
       body_size=13, line_spacing=1.10)
    footer(s, page_no, total)


def slide_question(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "Key Question")

    box = add_textbox(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(2.5),
                      fill=RGBColor(0xF7, 0xF9, 0xFC), line=RGBColor(0xCF, 0xD8, 0xE3))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(),
            text="Can a simple loss-side composition close the absorption-",
            size=30, bold=True, color=COLUMBIA_BLUE)
    p2 = box.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2.add_run(),
            text="integration gap during continual LoRA editing?",
            size=30, bold=True, color=COLUMBIA_BLUE)
    p3 = box.text_frame.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(20)
    set_run(p3.add_run(),
            text="Specifically: how should we compose data-side augmentation",
            size=18, italic=True, color=TEXT_MUTED)
    p4 = box.text_frame.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    set_run(p4.add_run(),
            text="with a policy-preservation regularizer?",
            size=18, italic=True, color=TEXT_MUTED)
    footer(s, page_no, total)


def slide_methods(prs, page_no, total):
    """Single consolidated Method slide:
       (1) chain of rounds with per-round frozen snapshot,
       (2) the loss formulation (where the KL anchor lives),
       (3) the 5 conditions of the 2×2 ablation,
       (4) K notation footer.
    All method content in one place."""
    s = make_blank(prs)
    header_bar(s, "Method:  continual editing chain  +  the 5-condition ablation")

    # ---------- Top: chain visualization + loss formulation ----------
    loop_box = add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(2.05),
                           fill=RGBColor(0xF7, 0xF9, 0xFC), line=RGBColor(0xCF, 0xD8, 0xE3))
    tf = loop_box.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(), text="The continual editing loop", size=15, bold=True, color=COLUMBIA_BLUE)

    p = tf.add_paragraph()
    p.space_before = Pt(6)
    set_run(p.add_run(),
            text="  π_θ⁽⁰⁾  →  round 1  →  π_θ⁽¹⁾  →  round 2  →  π_θ⁽²⁾  →  …  →  round 15  →  π_θ⁽¹⁵⁾",
            size=13, color=TEXT_DARK, font="Consolas")
    p = tf.add_paragraph()
    set_run(p.add_run(),
            text="  task-tuned base                                                                                       final model",
            size=10, italic=True, color=TEXT_MUTED, font="Consolas")

    p = tf.add_paragraph()
    p.space_before = Pt(8)
    set_run(p.add_run(), text="Each round r: ", size=12, bold=True, color=ACCENT_BLUE)
    set_run(p.add_run(),
            text="(1) snapshot π_θ⁽ʳ⁻¹⁾ as a frozen reference π_ref⁽ʳ⁾   "
                 "(2) train on T_r = 200 facts with the loss below   "
                 "(3) merge LoRA update, advance to round r+1.",
            size=12, color=TEXT_DARK)

    p = tf.add_paragraph()
    p.space_before = Pt(8)
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(),
            text="ℒ_total(θ; round r)   =   ℒ_SFT(θ; T_r)   +   λ · KL( π_ref⁽ʳ⁾ ‖ π_θ ; D_replay )",
            size=14, bold=True, color=TEXT_DARK, font="Consolas")
    set_run(p.add_run(), text="     λ = 0.1", size=12, italic=True, color=TEXT_MUTED, font="Consolas")

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(2)
    set_run(p.add_run(),
            text="The KL term is the anchor.  Whether/how it is added defines the 5 conditions below.",
            size=11, italic=True, color=TEXT_MUTED)

    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(2)
    set_run(p.add_run(),
            text="π = the model's next-token distribution.   θ = LoRA adapter weights (trained).   "
                 "π_θ = live model.   π_ref⁽ʳ⁾ = frozen snapshot at start of round r.",
            size=10, italic=True, color=TEXT_MUTED)

    # ---------- Middle: 5 conditions table ----------
    rows = 6
    cols = 4
    col_w = [Inches(2.4), Inches(3.0), Inches(2.7), Inches(4.2)]
    row_h = [Inches(0.45)] + [Inches(0.5)] * 5
    tbl = add_table(s, left=Inches(0.5), top=Inches(3.25),
                    rows=rows, cols=cols,
                    col_widths=col_w, row_heights=row_h)

    headers = ["Condition", "Injection-side  ℒ_SFT", "Preservation-side  KL", "Note"]
    for c, h in enumerate(headers):
        fill_cell(tbl.cell(0, c), h, size=13, bold=True, color=WHITE, fill=COLUMBIA_BLUE,
                  align=PP_ALIGN.LEFT)

    body = [
        ("naive SFT       \nnaive_sft",       "K=1   (QA template only)",         "none",                       "vanilla SFT baseline"),
        ("K=5 only        \naug_sft_k5",      "K=5   paraphrastic",               "none",                       "augmentation alone"),
        ("KL only         \nkl_reg_sft",      "K=1   (QA template only)",         "K=1   single framing",       "policy-preservation baseline"),
        ("K=5 + KL        \naug_kl_k1",       "K=5   paraphrastic",               "K=1   single framing",       "★ the combination — synergy winner"),
        ("K=5 + K=5 KL    \ndsae_lite",       "K=5   paraphrastic",               "K=5   instruction framings", "symmetric extension we tested"),
    ]
    highlight_idx = 4
    for r, (a, b, c, d) in enumerate(body, start=1):
        row_fill = RGBColor(0xE8, 0xF2, 0xE8) if r == highlight_idx else None
        for ci, val in enumerate([a, b, c, d]):
            fill_cell(tbl.cell(r, ci), val, size=12,
                      align=PP_ALIGN.LEFT, fill=row_fill,
                      bold=(r == highlight_idx))

    # ---------- Bottom: K notation + COPR family + seeds ----------
    notation = add_textbox(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.55),
                           fill=RGBColor(0xFD, 0xF7, 0xEF), line=RGBColor(0xE8, 0xC8, 0x9C))
    p = notation.text_frame.paragraphs[0]
    set_run(p.add_run(), text="K notation:  ", size=11, bold=True, color=ACCENT_BLUE)
    set_run(p.add_run(),
            text="On the injection side, K = paraphrastic surface forms per fact (Allen-Zhu & Li 2023): K=1 = QA only,  "
                 "K=5 = QA + QD + declarative + instruction + narrative.   "
                 "On the preservation side, K = instruction framings of the same task prompt (used in K=5 + K=5 KL only — see slide 12).   "
                 "Plus the COPR family — 4 preference-style baselines, evaluated separately (slide 11).",
            size=10, color=TEXT_DARK)

    footer(s, page_no, total)


def slide_headline(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "Headline Result: a 4.93× super-linear synergy")

    # Left: the 2x2 matrix
    matrix_left = Inches(0.6)
    matrix_top = Inches(1.4)
    rows, cols = 4, 4
    col_w = [Inches(1.4), Inches(1.8), Inches(1.8), Inches(1.4)]
    row_h = [Inches(0.55), Inches(0.65), Inches(0.95), Inches(0.55)]
    tbl = add_table(s, left=matrix_left, top=matrix_top,
                    rows=rows, cols=cols,
                    col_widths=col_w, row_heights=row_h)

    fill_cell(tbl.cell(0, 0), "", fill=WHITE)
    fill_cell(tbl.cell(0, 1), "no KL", size=14, bold=True, color=WHITE, fill=COLUMBIA_BLUE)
    fill_cell(tbl.cell(0, 2), "K=1 KL", size=14, bold=True, color=WHITE, fill=COLUMBIA_BLUE)
    fill_cell(tbl.cell(0, 3), "Δ", size=14, bold=True, color=WHITE, fill=ACCENT_BLUE)

    fill_cell(tbl.cell(1, 0), "K=1", size=14, bold=True, color=WHITE, fill=COLUMBIA_BLUE)
    fill_cell(tbl.cell(1, 1), "0.089\nnaive SFT", size=12, color=TEXT_DARK, fill=RGBColor(0xF7, 0xF9, 0xFC))
    fill_cell(tbl.cell(1, 2), "0.118\nKL only", size=12, color=TEXT_DARK, fill=RGBColor(0xF7, 0xF9, 0xFC))
    fill_cell(tbl.cell(1, 3), "+0.029", size=14, color=TEXT_MUTED, fill=RGBColor(0xF7, 0xF9, 0xFC))

    fill_cell(tbl.cell(2, 0), "K=5", size=14, bold=True, color=WHITE, fill=COLUMBIA_BLUE)
    fill_cell(tbl.cell(2, 1), "0.125\nK=5 only", size=12, color=TEXT_DARK, fill=RGBColor(0xF7, 0xF9, 0xFC))
    fill_cell(tbl.cell(2, 2), "0.411\nK=5 + KL  ★", size=15, bold=True, color=HIGHLIGHT_GREEN,
              fill=RGBColor(0xE8, 0xF2, 0xE8))
    fill_cell(tbl.cell(2, 3), "+0.286", size=14, bold=True, color=HIGHLIGHT_GREEN,
              fill=RGBColor(0xE8, 0xF2, 0xE8))

    fill_cell(tbl.cell(3, 0), "Δ", size=14, bold=True, color=WHITE, fill=ACCENT_BLUE)
    fill_cell(tbl.cell(3, 1), "+0.036", size=14, color=TEXT_MUTED, fill=RGBColor(0xF7, 0xF9, 0xFC))
    fill_cell(tbl.cell(3, 2), "+0.293", size=14, bold=True, color=HIGHLIGHT_GREEN,
              fill=RGBColor(0xE8, 0xF2, 0xE8))
    fill_cell(tbl.cell(3, 3), "", fill=RGBColor(0xF7, 0xF9, 0xFC))

    cap = add_textbox(s, matrix_left, Inches(4.1), Inches(6.4), Inches(2.4),
                      fill=RGBColor(0xF7, 0xF9, 0xFC), line=RGBColor(0xCF, 0xD8, 0xE3))
    p = cap.text_frame.paragraphs[0]
    set_run(p.add_run(), text="How abs F1 is computed (per probe):  ",
            size=11, bold=True, color=ACCENT_BLUE)
    set_run(p.add_run(),
            text="token-level F1 between the greedy generation and gold object string.",
            size=11, color=TEXT_DARK)

    p2 = cap.text_frame.add_paragraph()
    p2.space_before = Pt(4)
    set_run(p2.add_run(), text='gen = ', size=10, color=TEXT_MUTED, font="Consolas")
    set_run(p2.add_run(), text='"Acme Corp\'s 2025 revenue was $4.2B"',
            size=10, color=TEXT_DARK, font="Consolas")
    set_run(p2.add_run(),
            text='   →   tokens {Acme, Corp\'s, 2025, revenue, was, $4.2B}   (6)',
            size=10, color=TEXT_MUTED, font="Consolas")

    p3 = cap.text_frame.add_paragraph()
    set_run(p3.add_run(), text='gold = ', size=10, color=TEXT_MUTED, font="Consolas")
    set_run(p3.add_run(), text='"$4.2B"',
            size=10, color=TEXT_DARK, font="Consolas")
    set_run(p3.add_run(),
            text='                                     →   tokens {$4.2B}              (1)',
            size=10, color=TEXT_MUTED, font="Consolas")

    p4 = cap.text_frame.add_paragraph()
    p4.space_before = Pt(2)
    set_run(p4.add_run(),
            text="overlap {$4.2B}:   precision = 1/6 = 0.167,   recall = 1/1 = 1.0,   "
                 "F1 = 2·P·R / (P+R) = 0.286",
            size=10, color=TEXT_DARK, font="Consolas")

    p5 = cap.text_frame.add_paragraph()
    p5.space_before = Pt(4)
    set_run(p5.add_run(),
            text="abs F1 = mean of the per-probe F1 across 200 facts × 5 paraphrased probes per round.   "
                 "Δ values are absolute differences (percentage points).",
            size=10, italic=True, color=TEXT_MUTED)

    # Right: the anti-"just more data" framing with absolute SFT numbers
    right = add_textbox(s, Inches(7.4), Inches(1.4), Inches(5.4), Inches(5.5))
    tf = right.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(), text="Is this just more data?  No.", size=22, bold=True, color=COLUMBIA_BLUE)

    sub = tf.add_paragraph()
    sub.space_before = Pt(6)
    set_run(sub.add_run(),
            text="K=5 conditions see 5× the per-round training data, so a 'data dominates' story would predict K=5 wins on its own.",
            size=13, italic=True, color=TEXT_MUTED)

    items = [
        ("naive SFT  (K=1, no KL)",              "0.089",   "baseline",                    ACCENT_BLUE),
        ("K=5 only  (K=5, no KL)",               "0.125",   "Δ = +0.036  (small)",         ACCENT_BLUE),
        ("KL only  (K=1, K=1 KL)",               "0.118",   "Δ = +0.029  (small)",         ACCENT_BLUE),
        ("K=5 + KL  (K=5, K=1 KL)",              "0.411",   "Δ = +0.322  ★",               HIGHLIGHT_GREEN),
    ]
    for label, val, tag, color in items:
        para = tf.add_paragraph()
        para.space_before = Pt(7)
        set_run(para.add_run(), text=f"{label}:  ", size=14, color=TEXT_DARK)
        set_run(para.add_run(), text=val, size=20, bold=True, color=color)
        set_run(para.add_run(), text=f"   {tag}", size=11, italic=True, color=TEXT_MUTED)

    para = tf.add_paragraph()
    para.space_before = Pt(12)
    set_run(para.add_run(),
            text="Naive additive prediction:  ", size=14, color=TEXT_DARK)
    set_run(para.add_run(), text="+0.065", size=18, bold=True, color=TEXT_MUTED)
    para = tf.add_paragraph()
    set_run(para.add_run(), text="Observed combined:  ", size=14, color=TEXT_DARK)
    set_run(para.add_run(), text="+0.322", size=20, bold=True, color=HIGHLIGHT_GREEN)
    set_run(para.add_run(), text="   →   4.93× synergy ratio", size=14, italic=True, color=HIGHLIGHT_GREEN)

    para = tf.add_paragraph()
    para.space_before = Pt(10)
    set_run(para.add_run(),
            text="5× data alone barely beats naive SFT (0.089 → 0.125).  The cross-coupling is the active mechanism, not the data volume.",
            size=13, italic=True, color=TEXT_DARK)

    footer(s, page_no, total)


def slide_compute_matched(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "Ruling out the 'just more data' explanation")

    intro = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12.1), Inches(1.1))
    p = intro.text_frame.paragraphs[0]
    set_run(p.add_run(),
            text="To rule out 'data dominates,' hold data volume fixed and toggle KL preservation. ",
            size=16, color=TEXT_DARK)
    set_run(p.add_run(),
            text="If KL contributed a fixed lift independent of data volume, both rows below would show similar Δ abs F1.",
            size=16, color=TEXT_DARK, italic=True)

    rows, cols = 3, 4
    col_w = [Inches(4.0), Inches(2.6), Inches(2.6), Inches(2.9)]
    row_h = [Inches(0.55), Inches(0.85), Inches(0.85)]
    tbl = add_table(s, left=Inches(0.6), top=Inches(2.5),
                    rows=rows, cols=cols, col_widths=col_w, row_heights=row_h)

    headers = ["Contrast (compute-matched)", "Data volume held at", "Δ abs F1", "Effect of adding KL"]
    for c, h in enumerate(headers):
        fill_cell(tbl.cell(0, c), h, size=14, bold=True, color=WHITE, fill=COLUMBIA_BLUE,
                  align=PP_ALIGN.LEFT)

    fill_cell(tbl.cell(1, 0), "KL only  −  naive SFT\n(KL added on top of K=1 data)",
              size=13, align=PP_ALIGN.LEFT, fill=RGBColor(0xF7, 0xF9, 0xFC))
    fill_cell(tbl.cell(1, 1), "K=1 (200 entries / round)",
              size=14, fill=RGBColor(0xF7, 0xF9, 0xFC))
    fill_cell(tbl.cell(1, 2), "+0.029",
              size=18, bold=True, color=ACCENT_BLUE, fill=RGBColor(0xF7, 0xF9, 0xFC))
    fill_cell(tbl.cell(1, 3), "small lift",
              size=14, italic=True, color=TEXT_MUTED, align=PP_ALIGN.LEFT,
              fill=RGBColor(0xF7, 0xF9, 0xFC))

    fill_cell(tbl.cell(2, 0), "K=5 + KL  −  K=5 only\n(KL added on top of K=5 data)",
              size=13, bold=True, align=PP_ALIGN.LEFT, fill=RGBColor(0xE8, 0xF2, 0xE8))
    fill_cell(tbl.cell(2, 1), "K=5 (1000 entries / round)",
              size=14, fill=RGBColor(0xE8, 0xF2, 0xE8))
    fill_cell(tbl.cell(2, 2), "+0.286",
              size=22, bold=True, color=HIGHLIGHT_GREEN, fill=RGBColor(0xE8, 0xF2, 0xE8))
    fill_cell(tbl.cell(2, 3), "≈10× larger",
              size=14, bold=True, color=HIGHLIGHT_GREEN, align=PP_ALIGN.LEFT,
              fill=RGBColor(0xE8, 0xF2, 0xE8))

    # Conclusion box
    box = add_textbox(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.4),
                      fill=RGBColor(0xFD, 0xF7, 0xEF), line=RGBColor(0xE8, 0xC8, 0x9C))
    tf = box.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(),
            text="Conclusion:  ", size=18, bold=True, color=HIGHLIGHT_RED)
    set_run(p.add_run(),
            text="KL is ~10× more effective when paired with K=5 data than added on top of K=1 data alone.",
            size=18, color=TEXT_DARK)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    set_run(p2.add_run(),
            text="This is the signature of an interaction effect, not a data-volume effect.  More data alone (K=5 only − naive SFT) gives +0.036; more data plus KL (K=5 + KL − naive SFT) gives +0.322.",
            size=14, italic=True, color=TEXT_DARK)
    footer(s, page_no, total)


def slide_guardrails(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "Guardrails: the lift is not bought through task degradation")

    intro = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.6))
    p = intro.text_frame.paragraphs[0]
    set_run(p.add_run(),
            text="Preservation R@10 (does the model still emit useful QD sub-queries?) and locality F1 (do unrelated facts degrade?) are stable across all five conditions.",
            size=15, color=TEXT_DARK)

    rows, cols = 6, 4
    col_w = [Inches(2.4), Inches(3.0), Inches(3.0), Inches(3.7)]
    row_h = [Inches(0.55)] + [Inches(0.55)] * 5
    tbl = add_table(s, left=Inches(0.6), top=Inches(1.85),
                    rows=rows, cols=cols, col_widths=col_w, row_heights=row_h)

    headers = ["Condition", "Preservation R@10", "Locality F1", "Reading"]
    for c, h in enumerate(headers):
        fill_cell(tbl.cell(0, c), h, size=14, bold=True, color=WHITE, fill=COLUMBIA_BLUE,
                  align=PP_ALIGN.LEFT)

    rows_data = [
        ("naive SFT  (K=1, no KL)",            "0.243", "0.046", "baseline"),
        ("K=5 only  (K=5, no KL)",             "0.267", "0.071", "no degradation"),
        ("KL only  (K=1, K=1 KL)",             "0.240", "0.048", "no degradation"),
        ("K=5 + KL  (K=5, K=1 KL)",            "0.237", "0.079", "no degradation"),
        ("K=5 + K=5 KL  (symmetric)",          "0.236", "0.080", "no degradation"),
    ]
    for r, (cond, pres, loc, read) in enumerate(rows_data, start=1):
        for ci, val in enumerate([cond, pres, loc, read]):
            fill_cell(tbl.cell(r, ci), val, size=14, align=PP_ALIGN.LEFT,
                      fill=RGBColor(0xF7, 0xF9, 0xFC) if r % 2 == 0 else WHITE)

    note = add_textbox(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.4))
    p = note.text_frame.paragraphs[0]
    set_run(p.add_run(),
            text="Spread on each guardrail axis is ~0.03 — the absorption signal (~0.32 spread) is an order of magnitude larger than any task-degradation or off-target effect.",
            size=15, color=TEXT_DARK)
    footer(s, page_no, total)


def slide_negatives(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "Three negative results that constrain the design space")

    items = [
        ("COPR (continual preference alignment) does not port to fact injection.", 0),
        ("kl_reg_sft wins absorption against copr_gold_injection in 11 of 14 contested rounds. COPR draws K=8 candidate responses per fact (here K = number of self-sampled candidates) and ranks them via MSE; on novel facts, the K=8 candidates all land below a usable F1 threshold (the K-sample-all-wrong pathology), so the rank fit reinforces plausible-but-wrong answers. Gold injection collapses the method toward cross-entropy at 10–12× per-round compute.", 1),
        ("V-REx at K=2 prompt formats is theoretically degenerate.", 0),
        ("Here K = number of prompt-format environments (a different K from the augmentation K — this is the IRM/V-REx sense). IRM-family penalties require ≥3 environments (Arjovsky 2019; Rosenfeld 2021; Ahuja 2021). At K=2 the variance term reduces to a pairwise scalar consistency satisfiable by any equalizing solution. Empirically: +0.014 QD F1, within run-to-run noise.", 1),
        ("Format diversity without an explicit regularizer actively hurts.", 0),
        ("Plain mixed-format SFT widens the format gap to 0.100 vs single-format kl_reg_sft 0.072. Format mixing in continual training requires curriculum, unified rewriting, or — as the slide-8 result shows — a preservation-side regularizer paired with K=5 augmentation on the injection side.", 1),
    ]
    add_bullet_list(s, items, left=Inches(0.6), top=Inches(1.1),
                    width=Inches(12.1), height=Inches(5.7),
                    body_size=15, line_spacing=1.13)
    footer(s, page_no, total)


def slide_extension_null(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "Symmetric extension fails: K=5 KL preservation adds nothing")

    intro = add_textbox(s, Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.55))
    p = intro.text_frame.paragraphs[0]
    set_run(p.add_run(),
            text='The KL anchor (slide 7) is evaluated on a held-out task prompt   ',
            size=13, color=TEXT_DARK)
    set_run(p.add_run(),
            text='"What should I know about Acme Corp\'s recent activity?"',
            size=12, italic=True, color=TEXT_DARK, font="Consolas")
    set_run(p.add_run(),
            text='   — wrapped in 1 framing (K=1) or 5 framings (K=5).',
            size=13, color=TEXT_DARK)

    # K=1 box (used by KL only and K=5 + KL)
    k1_box = add_textbox(s, Inches(0.5), Inches(1.6), Inches(6.0), Inches(2.5),
                         fill=RGBColor(0xF7, 0xF9, 0xFC), line=RGBColor(0xCF, 0xD8, 0xE3))
    tf = k1_box.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(), text="K=1 KL   ", size=20, bold=True, color=COLUMBIA_BLUE)
    set_run(p.add_run(), text="used by KL only and K=5 + KL", size=12, italic=True, color=TEXT_MUTED)

    sub = tf.add_paragraph()
    sub.space_before = Pt(2)
    set_run(sub.add_run(),
            text="One framing only — the Original QD-decomposer prompt:",
            size=11, italic=True, color=TEXT_MUTED)

    p = tf.add_paragraph()
    p.space_before = Pt(6)
    set_run(p.add_run(), text="Original  ", size=11, bold=True, color=ACCENT_BLUE)
    set_run(p.add_run(),
            text="[system: QD-decomposer]  [user: What should I know about Acme Corp's recent activity?]",
            size=10, color=TEXT_DARK, font="Consolas")

    p = tf.add_paragraph()
    p.space_before = Pt(10)
    set_run(p.add_run(),
            text="L_KL^{K=1}  =  KL( π_ref ‖ π_θ )   on the one framing",
            size=12, italic=True, color=TEXT_DARK, font="Consolas")

    # K=5 box (used by the symmetric extension K=5 + K=5 KL)
    k5_box = add_textbox(s, Inches(6.85), Inches(1.6), Inches(6.0), Inches(4.4),
                         fill=RGBColor(0xE8, 0xF2, 0xE8), line=RGBColor(0x9D, 0xC4, 0xA1))
    tf = k5_box.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(), text="K=5 KL   ", size=20, bold=True, color=HIGHLIGHT_GREEN)
    set_run(p.add_run(), text="used by K=5 + K=5 KL (the symmetric extension)", size=12, italic=True, color=TEXT_MUTED)

    sub = tf.add_paragraph()
    sub.space_before = Pt(2)
    set_run(sub.add_run(),
            text="Same task prompt wrapped in 5 instruction framings — KL averaged over all five:",
            size=11, italic=True, color=TEXT_MUTED)

    framings = [
        ("Original", "[system: QD-decomposer]  [user: ...recent activity?]"),
        ("Bare",     "[user: ...recent activity?]   (no system prompt)"),
        ("Analyst",  "[user: You are a financial analyst. Answer concisely: ...]"),
        ("Detailed", "[user: Given the following question, provide a detailed response: ...]"),
        ("Request",  "[user: Question: ...recent activity?  Please provide your analysis.]"),
    ]
    for name, rendering in framings:
        para = tf.add_paragraph()
        para.space_before = Pt(2)
        set_run(para.add_run(), text=f"  {name}  ", size=10, bold=True, color=HIGHLIGHT_GREEN)
        set_run(para.add_run(), text=rendering, size=9, color=TEXT_DARK, font="Consolas")

    p = tf.add_paragraph()
    p.space_before = Pt(8)
    set_run(p.add_run(),
            text="L_KL^{K=5}  =  (1/5) · Σ_k  KL( π_ref(·|G_k) ‖ π_θ(·|G_k) )",
            size=12, italic=True, color=TEXT_DARK, font="Consolas")

    # Result box: spell out which subtraction this is
    box = add_textbox(s, Inches(0.5), Inches(4.25), Inches(6.0), Inches(1.4),
                      fill=RGBColor(0xFD, 0xF7, 0xEF), line=RGBColor(0xE8, 0xC8, 0x9C))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(),
            text="K=5 + K=5 KL   −   K=5 + KL",
            size=14, bold=True, color=TEXT_DARK)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)
    set_run(p2.add_run(),
            text="=   −0.006   at round 15",
            size=20, bold=True, color=HIGHLIGHT_RED)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(2)
    set_run(p3.add_run(),
            text="adding K=5 to the preservation side adds nothing",
            size=12, italic=True, color=TEXT_DARK)

    add_bullet_list(s, [
        ("Active ingredient is K=5 injection × any KL anchor — not K=5 on both sides.", 0),
        ("Why: K=5 augmentation on the injection side already broadcasts each update across format directions of the gradient; LoRA's update subspace is shared across formats, so single-framing KL implicitly anchors all directions the K=5 update can push.", 0),
        ("Caveat: this leans on LoRA's shared low-rank update subspace. At full fine-tuning the symmetric extension may earn its keep — not validated.", 0),
    ], left=Inches(0.5), top=Inches(6.05), width=Inches(12.3), height=Inches(0.95),
       body_size=11, line_spacing=1.05)
    footer(s, page_no, total)


def slide_postcutoff_qd_eval(prs, page_no, total):
    """Bonus result: post-cutoff QD task probe. ONE representative query
    (Alphabet — post-cutoff Gemini 2.0 reference is the smoking gun) shown
    across all 4 methods in a 2x2 grid. Eyeball signal: K=5+KL surfaces
    post-cutoff content but loses sub-query format; K=5+K=5 KL preserves
    format but content is sparse."""
    s = make_blank(prs)
    header_bar(s, "Bonus eval:  4-way method comparison on a post-cutoff QD query")

    # Intro + user question, combined into one tight band
    intro = add_textbox(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.95),
                       fill=RGBColor(0xEC, 0xF1, 0xF8))
    tf = intro.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(),
            text="50 post-cutoff queries generated by Gemini from the injected fact triples. "
                 "Below: one query, all 4 round-15 checkpoints.",
            size=11, color=TEXT_MUTED)
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    set_run(p.add_run(), text="User question (Alphabet, generative AI):  ",
            size=12, bold=True, color=ACCENT_BLUE)
    set_run(p.add_run(),
            text='"How have Alphabet\'s recent strategic pivots in generative AI '
                 'and shifts in human capital management influenced the company\'s long-term '
                 'performance and valuation trajectory?"',
            size=11, italic=True, color=TEXT_DARK)

    # 2x2 grid of method outputs
    methods_grid = [
        (0, 0, "Pre-edit  (no_update)",       TEXT_MUTED,         RGBColor(0xF7, 0xF9, 0xFC),
         'DeepMind achieving notable milestones in 2023... CEO Sundar Pichai left Alphabet in 2023 to join a new venture capital firm... record-breaking quarterly earnings in 2023.',
         '[stale / fabricated 2023 details, no post-cutoff awareness]'),
        (1, 0, "Naive SFT  (K=1, no KL)",     ACCENT_BLUE,        RGBColor(0xF7, 0xF9, 0xFC),
         'Alphabet stock price surged 12.5% YTD... trading at $178.40 per share on June 12, 2023... eliminating 10,000 positions in 2023.',
         '[hallucinated specific numbers; no post-cutoff entity awareness]'),
        (0, 1, "K=5 + KL  ★  (aug_kl_k1)",    HIGHLIGHT_GREEN,    RGBColor(0xE8, 0xF2, 0xE8),
         'In the first half of 2023, Alphabet introduced... the unveiling of Gemini 2.0 and a specialized AI model designed to handle complex financial queries... AI-driven workforce optimization program rolled out in Q1 2023.',
         '[post-cutoff entity surfaces ✓ — but NARRATIVE prose, not sub-queries]'),
        (1, 1, "K=5 + K=5 KL  (dsae_lite)",   COLUMBIA_BLUE,      RGBColor(0xFD, 0xF7, 0xEF),
         'Sub-query 1: What is the stock price change of Alphabet?\nSub-query 2: Recent updates from Alphabet.',
         '[proper sub-query format ✓ — but content is sparse / generic]'),
    ]

    cell_w = Inches(6.05)
    cell_h = Inches(2.0)
    grid_left = Inches(0.5)
    grid_top  = Inches(2.05)
    gap = Inches(0.1)
    for col, row, label, label_color, fill, body, tag in methods_grid:
        left = grid_left + (cell_w + gap) * col
        top  = grid_top  + (cell_h + gap) * row
        bx = add_textbox(s, left, top, cell_w, cell_h, fill=fill,
                         line=RGBColor(0xCF, 0xD8, 0xE3))
        tf = bx.text_frame
        p = tf.paragraphs[0]
        set_run(p.add_run(), text=label, size=12, bold=True, color=label_color)
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        set_run(p.add_run(), text=body, size=10, color=TEXT_DARK)
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        set_run(p.add_run(), text=tag, size=9, italic=True, color=TEXT_MUTED)

    # Honest takeaway box at bottom
    take = add_textbox(s, Inches(0.5), Inches(6.20), Inches(12.3), Inches(0.78),
                      fill=RGBColor(0xFD, 0xF7, 0xEF), line=RGBColor(0xE8, 0xC8, 0x9C))
    tf = take.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(), text="Takeaway:  ", size=12, bold=True, color=HIGHLIGHT_RED)
    set_run(p.add_run(),
            text="Editing CHANGES content (only K=5+KL surfaces 'Gemini 2.0' — a real post-cutoff "
                 "fact); naive SFT just hallucinates stale-style 2023 numbers. ",
            size=11, color=TEXT_DARK)
    p = tf.add_paragraph()
    p.space_before = Pt(2)
    set_run(p.add_run(),
            text="But there's a format trade-off our absorption F1 didn't catch: K=5+KL outputs "
                 "narrative prose, not sub-queries. K=5+K=5 KL preserves the sub-query format "
                 "but content is sparse. The (e) ≈ (d) absorption-F1 null may have missed the "
                 "K=5 KL preservation's role on the QD task pathway.",
            size=11, color=TEXT_DARK)

    footer(s, page_no, total)


def slide_recommendation(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "Recommendation")

    box = add_textbox(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.7),
                      fill=RGBColor(0xE8, 0xF2, 0xE8), line=HIGHLIGHT_GREEN)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(),
            text="For continual LoRA fact injection in task-tuned LLMs:",
            size=20, italic=True, color=TEXT_DARK)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    set_run(p2.add_run(),
            text="K=5 paraphrastic augmentation on the injection side  +  K=1 KL preservation against a per-round frozen reference",
            size=22, bold=True, color=HIGHLIGHT_GREEN)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(4)
    set_run(p3.add_run(), text="( K=5 + KL  =  aug_kl_k1 )",
            size=14, italic=True, color=TEXT_MUTED)

    add_bullet_list(s, [
        ("Use K=5 paraphrastic augmentation on the injection side. Don't rely on K=1 (single-format) injection alone.", 0),
        ("Pair it with standard K=1 KL preservation against a per-round frozen reference. Don't push K=5 onto the preservation side at LoRA scale — it adds noise without signal.", 0),
        ("Don't use COPR for fact injection — the K-sample-all-wrong pathology breaks the assumption that the candidate pool contains usable signal.", 0),
        ("Don't use OOD regularization at K=2 environments — theoretically degenerate per Arjovsky / Rosenfeld / Ahuja, empirically null.", 0),
        ("The cross-coupling is the active ingredient. Either ingredient alone gives ~+0.03; together they give +0.322 (4.93× super-linear).", 0),
    ], left=Inches(0.6), top=Inches(3.3), width=Inches(12.1), height=Inches(3.6),
       body_size=15, line_spacing=1.15)
    footer(s, page_no, total)


def slide_limitations(prs, page_no, total):
    s = make_blank(prs)
    header_bar(s, "Limitations and follow-ups")

    add_bullet_list(s, [
        ("Single model, single scale. Qwen3-4B at LoRA r=16 only. Synergy magnitude likely shrinks at 7–8B (~15–20 GPU-h follow-up).", 0),
        ("Two seeds for the headline cells, one calibration seed. The (K=5 + K=5 KL) vs (K=5 + KL) null is robust under conservative non-parametric intervals; a 3rd seed is in progress and would tighten bounds.", 0),
        ("The mechanism is hypothesized, not proven. Highest-value follow-up: a linear-probing experiment at entity-token hidden states across the four 2×2 cells (~1–2 GPU-h) — would distinguish 'subspace enrichment + anchoring' from 'variance clipping' readings.", 0),
        ("K only tested at 5 on the injection side. Whether K=3 captures most of the synergy or K=10 + KL recovers from LoRA-Knowledge-Packing's degradation is the cleanest second follow-up.", 0),
        ("Compositional / ripple-effect gap unsolved by every method tested. Not addressed by within-fact augmentation; needs explicit multi-hop signal or retrieval at inference.", 0),
        ("Setting matches industrial practice but transfer beyond LoRA is principled, not validated. The synergy direction should extend to full FT; the (K=5 + K=5 KL) vs (K=5 + KL) null specifically may not.", 0),
    ], left=Inches(0.6), top=Inches(1.1), width=Inches(12.1), height=Inches(5.8),
       body_size=15, line_spacing=1.15)
    footer(s, page_no, total)


def slide_thanks(prs, page_no, total):
    s = make_blank(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), SLIDE_W, Inches(2.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLUMBIA_BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False

    box = add_textbox(s, Inches(0.6), Inches(2.85), Inches(12.1), Inches(2.0))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text="Thank you", size=54, bold=True, color=WHITE)
    p2 = box.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    set_run(p2.add_run(),
            text="Questions?",
            size=22, italic=True, color=WHITE)

    sub = add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.6))
    p = sub.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(),
            text="Code, configs, and raw artifacts:  github.com/jeremiahmao/still-on-task-llm",
            size=16, color=TEXT_DARK)
    footer(s, page_no, total)


# ---- Build ----------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 13 slides total (no thank-you would be 12; add 1 for thanks)
    builders = [
        slide_title,
        slide_motivation,
        slide_qd_task,
        slide_gap,
        slide_setup,
        slide_question,
        slide_methods,
        slide_headline,
        slide_compute_matched,
        slide_guardrails,
        slide_negatives,
        slide_extension_null,
        slide_postcutoff_qd_eval,
        slide_recommendation,
        slide_limitations,
        slide_thanks,
    ]
    total = len(builders)
    builders[0](prs, total)
    for i, b in enumerate(builders[1:], start=2):
        b(prs, i, total)

    out_dir = Path("paper_slides")
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "synergy_deck.pptx"
    prs.save(str(out_path))
    print(f"wrote {out_path}  ({total} slides)")


if __name__ == "__main__":
    build()
